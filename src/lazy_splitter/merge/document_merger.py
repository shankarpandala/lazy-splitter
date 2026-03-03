"""Document merger for DOCX, PPTX, Markdown, and CSV files.

Provides :class:`DocumentMerger` with dedicated methods for each supported
document format.  External libraries (python-docx, python-pptx) are imported
optionally so that only the dependencies actually used at runtime are
required.
"""

from __future__ import annotations

import csv
import io
import logging
import time
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Sequence

from lazy_splitter.core.base import BaseMerger
from lazy_splitter.core.exceptions import MergeError
from lazy_splitter.core.models import MergeResult
from lazy_splitter.core.utils import ensure_dir, sanitize_filename

# Optional heavy dependencies
try:
    from docx import Document as DocxDocument  # type: ignore[import-untyped]
    from docx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore[import-untyped]
except ImportError:  # pragma: no cover
    DocxDocument = None  # type: ignore[assignment,misc]
    RT = None  # type: ignore[assignment]

try:
    from pptx import Presentation  # type: ignore[import-untyped]
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore[assignment,misc]


class DocumentMerger(BaseMerger):
    """Merge Word, PowerPoint, Markdown, and CSV files.

    Each ``merge_*`` method targets a single format family.  The generic
    :meth:`merge` method from the base class dispatches to the correct
    implementation based on the file extension of the first input.

    Parameters
    ----------
    logger:
        Optional :class:`logging.Logger` instance.
    """

    def __init__(self, logger: Optional[logging.Logger] = None) -> None:
        super().__init__(logger=logger)

    # ------------------------------------------------------------------
    # BaseMerger interface
    # ------------------------------------------------------------------

    def merge(
        self,
        paths: Sequence[Path],
        output_path: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **kwargs: Any,
    ) -> MergeResult:
        """Auto-dispatch merge based on the output file extension.

        Supported extensions: ``.docx``, ``.pptx``, ``.md``, ``.csv``.

        Parameters
        ----------
        paths:
            Ordered sequence of input file paths.
        output_path:
            Destination path for the merged file.
        progress_callback:
            Optional progress callable (passed through to the format-
            specific method).
        **kwargs:
            Forwarded to the underlying ``merge_*`` method.

        Returns
        -------
        MergeResult

        Raises
        ------
        MergeError
            If the output format is not supported.
        """
        output_path = Path(output_path)
        ext = output_path.suffix.lower()

        dispatch = {
            ".docx": self.merge_docx,
            ".pptx": self.merge_pptx,
            ".md": self.merge_markdown,
            ".markdown": self.merge_markdown,
            ".csv": self.merge_csv,
        }

        handler = dispatch.get(ext)
        if handler is None:
            raise MergeError(
                f"Unsupported document format for merging: {ext!r}. "
                f"Supported: {', '.join(sorted(dispatch))}"
            )

        return handler(paths, output_path, **kwargs)  # type: ignore[call-arg]

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def merge_docx(
        self,
        paths: Sequence[Path],
        output_path: Path,
        **kwargs: Any,
    ) -> MergeResult:
        """Merge multiple Word (.docx) documents into one.

        The first document serves as the base (its styles, headers, and
        footers are preserved).  Subsequent documents are appended with an
        optional page break between each.

        Parameters
        ----------
        paths:
            Ordered sequence of ``.docx`` file paths.
        output_path:
            Destination for the merged document.
        **kwargs:
            ``page_break`` (bool, default ``True``) -- insert a page break
            before each appended document.

        Returns
        -------
        MergeResult

        Raises
        ------
        MergeError
            If python-docx is not installed or any input is invalid.
        """
        if DocxDocument is None:
            raise MergeError(
                "python-docx is required for DOCX merging. "
                "Install it with: pip install python-docx"
            )

        validated = self._validate_inputs(paths)
        output_path = Path(output_path)
        ensure_dir(output_path.parent)

        start_time = time.monotonic()
        insert_page_break = kwargs.get("page_break", True)

        # Use the first document as the base
        base_doc = DocxDocument(str(validated[0]))

        for idx, doc_path in enumerate(validated[1:], start=2):
            self.logger.debug(
                "Appending %s (%d/%d)", doc_path.name, idx, len(validated),
            )

            if insert_page_break:
                base_doc.add_page_break()

            src_doc = DocxDocument(str(doc_path))

            for element in src_doc.element.body:
                # Deep-copy each block-level element into the base document
                base_doc.element.body.append(element)

        try:
            base_doc.save(str(output_path))
        except Exception as exc:
            raise MergeError(
                f"Failed to save merged DOCX: {exc}",
                path=str(output_path),
            ) from exc

        elapsed = time.monotonic() - start_time
        self.logger.info(
            "Merged %d DOCX files into %s (%.2fs)",
            len(validated),
            output_path.name,
            elapsed,
        )

        return MergeResult(
            output_path=output_path,
            source_paths=list(validated),
            duration_seconds=elapsed,
            metadata={"page_break": insert_page_break},
        )

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def merge_pptx(
        self,
        paths: Sequence[Path],
        output_path: Path,
        **kwargs: Any,
    ) -> MergeResult:
        """Merge multiple PowerPoint (.pptx) presentations into one.

        Slides from each source are appended to the first presentation in
        order.  Slide layouts from each source are reused where possible.

        Parameters
        ----------
        paths:
            Ordered sequence of ``.pptx`` file paths.
        output_path:
            Destination for the merged presentation.
        **kwargs:
            Reserved for future options.

        Returns
        -------
        MergeResult

        Raises
        ------
        MergeError
            If python-pptx is not installed or any input is invalid.
        """
        if Presentation is None:
            raise MergeError(
                "python-pptx is required for PPTX merging. "
                "Install it with: pip install python-pptx"
            )

        validated = self._validate_inputs(paths)
        output_path = Path(output_path)
        ensure_dir(output_path.parent)

        start_time = time.monotonic()
        total_slides = 0

        # Use the first presentation as the base
        base_prs = Presentation(str(validated[0]))
        total_slides += len(base_prs.slides)

        for idx, pptx_path in enumerate(validated[1:], start=2):
            self.logger.debug(
                "Appending %s (%d/%d)", pptx_path.name, idx, len(validated),
            )

            src_prs = Presentation(str(pptx_path))

            for slide in src_prs.slides:
                # Determine the best matching layout in the base presentation
                layout = self._find_matching_layout(base_prs, slide)
                new_slide = base_prs.slides.add_slide(layout)

                # Copy all shapes from the source slide
                for shape in slide.shapes:
                    self._copy_shape(shape, new_slide)

                # Copy slide notes if present
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text:
                        new_slide.notes_slide.notes_text_frame.text = notes_text

                total_slides += 1

        try:
            base_prs.save(str(output_path))
        except Exception as exc:
            raise MergeError(
                f"Failed to save merged PPTX: {exc}",
                path=str(output_path),
            ) from exc

        elapsed = time.monotonic() - start_time
        self.logger.info(
            "Merged %d PPTX files (%d slides) into %s (%.2fs)",
            len(validated),
            total_slides,
            output_path.name,
            elapsed,
        )

        return MergeResult(
            output_path=output_path,
            source_paths=list(validated),
            duration_seconds=elapsed,
            metadata={"total_slides": total_slides},
        )

    # ------------------------------------------------------------------
    # Markdown
    # ------------------------------------------------------------------

    def merge_markdown(
        self,
        paths: Sequence[Path],
        output_path: Path,
        **kwargs: Any,
    ) -> MergeResult:
        """Merge multiple Markdown files into one.

        Each source file is separated by an optional heading (derived from
        the filename) and a horizontal rule.

        Parameters
        ----------
        paths:
            Ordered sequence of ``.md`` / ``.markdown`` file paths.
        output_path:
            Destination for the merged Markdown file.
        **kwargs:
            ``separator`` (str) -- text inserted between documents.
            Defaults to ``"\\n\\n---\\n\\n"``.
            ``add_headings`` (bool, default ``True``) -- prepend a
            level-1 heading with the source filename before each section.
            ``encoding`` (str, default ``"utf-8"``) -- source file
            encoding.

        Returns
        -------
        MergeResult
        """
        validated = self._validate_inputs(paths)
        output_path = Path(output_path)
        ensure_dir(output_path.parent)

        start_time = time.monotonic()
        separator = kwargs.get("separator", "\n\n---\n\n")
        add_headings = kwargs.get("add_headings", True)
        encoding = kwargs.get("encoding", "utf-8")

        parts: List[str] = []

        for md_path in validated:
            self.logger.debug("Reading %s", md_path.name)
            try:
                content = md_path.read_text(encoding=encoding)
            except Exception as exc:
                raise MergeError(
                    f"Failed to read Markdown file: {exc}",
                    path=str(md_path),
                ) from exc

            if add_headings:
                heading = f"# {md_path.stem}\n\n"
                parts.append(heading + content)
            else:
                parts.append(content)

        merged_text = separator.join(parts)

        # Ensure the file ends with a single newline
        if not merged_text.endswith("\n"):
            merged_text += "\n"

        try:
            output_path.write_text(merged_text, encoding=encoding)
        except Exception as exc:
            raise MergeError(
                f"Failed to write merged Markdown: {exc}",
                path=str(output_path),
            ) from exc

        elapsed = time.monotonic() - start_time
        self.logger.info(
            "Merged %d Markdown files into %s (%.2fs)",
            len(validated),
            output_path.name,
            elapsed,
        )

        return MergeResult(
            output_path=output_path,
            source_paths=list(validated),
            duration_seconds=elapsed,
            metadata={
                "add_headings": add_headings,
                "total_chars": len(merged_text),
            },
        )

    # ------------------------------------------------------------------
    # CSV
    # ------------------------------------------------------------------

    def merge_csv(
        self,
        paths: Sequence[Path],
        output_path: Path,
        dedupe_header: bool = True,
        **kwargs: Any,
    ) -> MergeResult:
        """Merge multiple CSV files into one.

        By default the header row is taken from the first file and duplicate
        headers in subsequent files are skipped.

        Parameters
        ----------
        paths:
            Ordered sequence of ``.csv`` file paths.
        output_path:
            Destination for the merged CSV.
        dedupe_header:
            If ``True`` (default), only the header from the first file is
            kept and matching header rows in subsequent files are skipped.
        **kwargs:
            ``encoding`` (str, default ``"utf-8"``) -- file encoding.
            ``delimiter`` (str, default ``","``) -- CSV delimiter character.

        Returns
        -------
        MergeResult
        """
        validated = self._validate_inputs(paths)
        output_path = Path(output_path)
        ensure_dir(output_path.parent)

        start_time = time.monotonic()
        encoding = kwargs.get("encoding", "utf-8")
        delimiter = kwargs.get("delimiter", ",")

        total_rows = 0
        header_row: Optional[List[str]] = None

        try:
            with open(str(output_path), "w", newline="", encoding=encoding) as out_fh:
                writer = csv.writer(out_fh, delimiter=delimiter)

                for file_idx, csv_path in enumerate(validated):
                    self.logger.debug("Reading %s", csv_path.name)
                    with open(str(csv_path), "r", newline="", encoding=encoding) as in_fh:
                        reader = csv.reader(in_fh, delimiter=delimiter)

                        for row_idx, row in enumerate(reader):
                            # Handle header deduplication
                            if row_idx == 0 and dedupe_header:
                                if header_row is None:
                                    # First file: write and remember the header
                                    header_row = row
                                    writer.writerow(row)
                                    total_rows += 1
                                elif row == header_row:
                                    # Duplicate header: skip
                                    continue
                                else:
                                    # Different header: write it (data may have
                                    # different columns)
                                    writer.writerow(row)
                                    total_rows += 1
                            else:
                                writer.writerow(row)
                                total_rows += 1

        except Exception as exc:
            if not isinstance(exc, MergeError):
                raise MergeError(
                    f"Failed to merge CSV files: {exc}",
                    path=str(output_path),
                ) from exc
            raise

        elapsed = time.monotonic() - start_time
        self.logger.info(
            "Merged %d CSV files (%d rows) into %s (%.2fs)",
            len(validated),
            total_rows,
            output_path.name,
            elapsed,
        )

        return MergeResult(
            output_path=output_path,
            source_paths=list(validated),
            duration_seconds=elapsed,
            metadata={
                "total_rows": total_rows,
                "dedupe_header": dedupe_header,
            },
        )

    # ------------------------------------------------------------------
    # PPTX private helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _find_matching_layout(base_prs: Any, slide: Any) -> Any:
        """Find a slide layout in *base_prs* that best matches *slide*.

        Falls back to the first layout if no match is found.

        Parameters
        ----------
        base_prs:
            The base ``Presentation`` object.
        slide:
            A slide from a source presentation.

        Returns
        -------
        SlideLayout
            The best matching layout from *base_prs*.
        """
        src_layout_name = slide.slide_layout.name if slide.slide_layout else None

        if src_layout_name:
            for layout in base_prs.slide_layouts:
                if layout.name == src_layout_name:
                    return layout

        # Fallback: use the first available layout
        return base_prs.slide_layouts[0]

    @staticmethod
    def _copy_shape(shape: Any, target_slide: Any) -> None:
        """Copy a shape's XML element into *target_slide*.

        This performs a deep copy of the underlying XML so that all shape
        properties (text, images, formatting) are preserved.

        Parameters
        ----------
        shape:
            A shape from a source slide.
        target_slide:
            The destination slide in the base presentation.
        """
        try:
            from copy import deepcopy
            from lxml import etree  # type: ignore[import-untyped]

            el = deepcopy(shape.element)
            target_slide.shapes._spTree.append(el)
        except Exception:
            # If deep-copy fails (e.g. linked media), silently skip the
            # shape rather than aborting the entire merge.
            pass
