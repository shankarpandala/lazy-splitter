"""Document section detection for DOCX, PPTX, XLSX, Markdown, LaTeX, HTML, CSV, JSON, and XML.

The :class:`DocumentDetector` auto-dispatches to a format-specific handler
based on the file extension, producing a :class:`~lazy_splitter.core.models.DetectionResult`
whose *segments* list contains :class:`~lazy_splitter.document.models.DocumentSection`
instances.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
from pathlib import Path
from typing import Any, Dict, List, Optional
from xml.etree import ElementTree

from lazy_splitter.core.base import BaseDetector
from lazy_splitter.core.exceptions import DetectionError, FileTypeError
from lazy_splitter.core.models import DetectionResult
from lazy_splitter.document.models import (
    DocumentInfo,
    DocumentSection,
    FileType,
    SectionType,
    file_type_from_path,
)

# ---------------------------------------------------------------------------
# Optional heavyweight imports -- gracefully degrade when not installed.
# ---------------------------------------------------------------------------
try:
    from docx import Document as DocxDocument  # type: ignore[import-untyped]
    from docx.opc.exceptions import PackageNotFoundError as DocxPackageError  # type: ignore[import-untyped]

    _HAS_DOCX = True
except ImportError:  # pragma: no cover
    _HAS_DOCX = False

try:
    from pptx import Presentation  # type: ignore[import-untyped]

    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _HAS_PPTX = False

try:
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    _HAS_OPENPYXL = True
except ImportError:  # pragma: no cover
    _HAS_OPENPYXL = False


class DocumentDetector(BaseDetector):
    """Detect logical sections within a variety of document formats.

    The public entry-point is :meth:`detect`, which inspects the file
    extension and delegates to one of the private ``_detect_*`` helpers.

    Parameters
    ----------
    logger:
        Optional logger instance.  Falls back to a module-level logger.
    """

    # Heading styles used by python-docx (English locale, case-insensitive prefix).
    _DOCX_HEADING_RE = re.compile(r"^Heading\s*(\d+)$", re.IGNORECASE)

    # LaTeX sectioning commands in descending order of depth.
    _LATEX_SECTION_RE = re.compile(
        r"\\(chapter|section|subsection|subsubsection)\s*(?:\[.*?\])?\s*\{([^}]*)\}"
    )
    _LATEX_LEVEL_MAP: Dict[str, int] = {
        "chapter": 1,
        "section": 2,
        "subsection": 3,
        "subsubsection": 4,
    }

    # Markdown ATX heading (lines starting with 1-6 ``#``).
    _MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)(?:\s+#+\s*)?$", re.MULTILINE)

    # HTML heading tags (h1 .. h6).
    _HTML_HEADING_RE = re.compile(
        r"<(h[1-6])\b[^>]*>(.*?)</\1>", re.IGNORECASE | re.DOTALL
    )

    # ------------------------------------------------------------------ #
    # Public API                                                          #
    # ------------------------------------------------------------------ #

    def detect(
        self,
        input_path: Path,
        strategy: str = "auto",
        **kwargs: Any,
    ) -> DetectionResult:
        """Detect sections in the document at *input_path*.

        Parameters
        ----------
        input_path:
            Path to the document file.
        strategy:
            Detection strategy.  ``"auto"`` inspects the file extension and
            picks the appropriate handler.  Pass a format name (e.g.
            ``"markdown"``, ``"docx"``) to force a specific handler.
        **kwargs:
            Strategy-specific options forwarded to the handler:

            * *heading_level* (int) -- maximum heading level for DOCX /
              Markdown (default ``3``).
            * *rows_per_chunk* (int) -- rows per chunk for CSV detection
              (default ``1000``).

        Returns
        -------
        DetectionResult

        Raises
        ------
        FileTypeError
            If the file type cannot be determined or is unsupported.
        DetectionError
            If detection fails for any reason.
        """
        path = Path(input_path)
        if not path.is_file():
            raise DetectionError(
                f"File not found: {path}", path=str(path)
            )

        file_type = file_type_from_path(path)
        if strategy != "auto":
            # Allow callers to force a handler by name.
            try:
                file_type = FileType(strategy)
            except ValueError:
                raise FileTypeError(
                    f"Unknown strategy / file type: {strategy!r}",
                    strategy=strategy,
                )

        if file_type == FileType.UNKNOWN:
            raise FileTypeError(
                f"Unsupported file type: {path.suffix!r}", path=str(path)
            )

        dispatch = {
            FileType.DOCX: self._detect_docx,
            FileType.PPTX: self._detect_pptx,
            FileType.XLSX: self._detect_xlsx,
            FileType.MARKDOWN: self._detect_markdown,
            FileType.LATEX: self._detect_latex,
            FileType.HTML: self._detect_html,
            FileType.CSV: self._detect_csv,
            FileType.JSON: self._detect_json,
            FileType.XML: self._detect_xml,
        }

        handler = dispatch.get(file_type)
        if handler is None:
            raise FileTypeError(
                f"No handler for file type {file_type.value!r}",
                file_type=file_type.value,
            )

        try:
            sections = handler(path, **kwargs)
        except (DetectionError, FileTypeError):
            raise
        except Exception as exc:
            raise DetectionError(
                f"Detection failed for {path.name}: {exc}",
                path=str(path),
                original_error=str(exc),
            ) from exc

        self.logger.info(
            "Detected %d section(s) in %s [%s]",
            len(sections),
            path.name,
            file_type.value,
        )

        return DetectionResult(
            chapters=sections,
            strategy_used=file_type.value,
            total_items=len(sections),
            metadata={"file_type": file_type.value},
            source_path=str(path),
            file_type=file_type.value,
        )

    # ------------------------------------------------------------------ #
    # Convenience                                                         #
    # ------------------------------------------------------------------ #

    def get_document_info(self, path: Path) -> DocumentInfo:
        """Return high-level metadata for the document at *path*.

        Parameters
        ----------
        path:
            Path to the document file.

        Returns
        -------
        DocumentInfo
        """
        path = Path(path)
        file_type = file_type_from_path(path)
        info = DocumentInfo(path=path, file_type=file_type)

        try:
            result = self.detect(path)
            info.section_count = result.chapter_count
        except Exception:  # noqa: BLE001
            pass

        # Format-specific enrichment.
        try:
            if file_type == FileType.DOCX and _HAS_DOCX:
                self._enrich_docx_info(path, info)
            elif file_type == FileType.PPTX and _HAS_PPTX:
                self._enrich_pptx_info(path, info)
            elif file_type == FileType.XLSX and _HAS_OPENPYXL:
                self._enrich_xlsx_info(path, info)
            elif file_type in (
                FileType.MARKDOWN,
                FileType.LATEX,
                FileType.HTML,
                FileType.CSV,
                FileType.JSON,
                FileType.XML,
            ):
                self._enrich_text_info(path, info)
        except Exception:  # noqa: BLE001
            pass

        return info

    # ------------------------------------------------------------------ #
    # DOCX                                                                #
    # ------------------------------------------------------------------ #

    def _detect_docx(
        self,
        path: Path,
        heading_level: int = 3,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect sections in a Word document by heading styles.

        Parameters
        ----------
        path:
            Path to the ``.docx`` file.
        heading_level:
            Maximum heading level to consider (e.g. ``2`` detects Heading 1
            and Heading 2 only).

        Returns
        -------
        List[DocumentSection]
        """
        if not _HAS_DOCX:
            raise DetectionError(
                "python-docx is required for DOCX detection.  "
                "Install it with: pip install python-docx",
                path=str(path),
            )

        try:
            doc = DocxDocument(str(path))
        except Exception as exc:
            raise DetectionError(
                f"Failed to open DOCX file: {exc}", path=str(path)
            ) from exc

        sections: List[DocumentSection] = []
        current_offset = 0

        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""
            match = self._DOCX_HEADING_RE.match(style_name)
            if match:
                level = int(match.group(1))
                if level <= heading_level:
                    title = para.text.strip() or f"Heading {level} (untitled)"
                    preview = ""
                    sections.append(
                        DocumentSection(
                            title=title,
                            start_index=current_offset,
                            end_index=current_offset,  # updated below
                            level=level,
                            section_type=SectionType.HEADING,
                            detection_method="docx_heading_style",
                            confidence=1.0,
                            content_preview=preview,
                        )
                    )
            current_offset += len(para.text) + 1  # +1 for paragraph break

        # Back-fill end indices so each section runs until the next one starts.
        total_length = current_offset
        for i, sec in enumerate(sections):
            if i + 1 < len(sections):
                sec.end_index = sections[i + 1].start_index
            else:
                sec.end_index = total_length

        return sections

    # ------------------------------------------------------------------ #
    # PPTX                                                                #
    # ------------------------------------------------------------------ #

    def _detect_pptx(
        self,
        path: Path,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect slides (and sections if present) in a PowerPoint file.

        Parameters
        ----------
        path:
            Path to the ``.pptx`` file.

        Returns
        -------
        List[DocumentSection]
        """
        if not _HAS_PPTX:
            raise DetectionError(
                "python-pptx is required for PPTX detection.  "
                "Install it with: pip install python-pptx",
                path=str(path),
            )

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise DetectionError(
                f"Failed to open PPTX file: {exc}", path=str(path)
            ) from exc

        sections: List[DocumentSection] = []

        # Check for native sections (PPTX section objects are not directly
        # exposed by python-pptx, so we detect them via the slide XML when
        # available).  Fall back to one-section-per-slide.
        section_starts = self._find_pptx_sections(prs)

        if section_starts:
            slides = list(prs.slides)
            for idx, (start_slide, sec_title) in enumerate(section_starts):
                end_slide = (
                    section_starts[idx + 1][0]
                    if idx + 1 < len(section_starts)
                    else len(slides)
                )
                sections.append(
                    DocumentSection(
                        title=sec_title or f"Section {idx + 1}",
                        start_index=start_slide,
                        end_index=end_slide,
                        level=1,
                        section_type=SectionType.SECTION,
                        detection_method="pptx_section",
                        confidence=1.0,
                        content_preview="",
                    )
                )
        else:
            # One section per slide.
            for idx, slide in enumerate(prs.slides):
                title = self._extract_slide_title(slide) or f"Slide {idx + 1}"
                sections.append(
                    DocumentSection(
                        title=title,
                        start_index=idx,
                        end_index=idx + 1,
                        level=1,
                        section_type=SectionType.SLIDE,
                        detection_method="pptx_slide",
                        confidence=0.9,
                        content_preview=title,
                    )
                )

        return sections

    @staticmethod
    def _extract_slide_title(slide: Any) -> str:
        """Return the title text of a slide, or ``""`` if none found."""
        if slide.shapes.title is not None:
            return slide.shapes.title.text.strip()
        # Fallback: look for any text frame with a title placeholder.
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    return text
        return ""

    @staticmethod
    def _find_pptx_sections(prs: Any) -> List[Any]:
        """Attempt to read native PPTX sections from the presentation XML.

        Returns a list of ``(start_slide_index, section_title)`` tuples, or
        an empty list when the presentation has no section information.
        """
        try:
            # python-pptx exposes sections via the presentation part XML
            # (Office 2019+).  The ``<p:sectionLst>`` element lives inside
            # ``<p:presentation>``.
            prs_elem = prs.presentation.element  # type: ignore[union-attr]
            nsmap = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
            }
            section_list = prs_elem.find(".//p14:sectionLst", nsmap)
            if section_list is None:
                return []

            results = []
            slide_ids = [
                slide.slide_id for slide in prs.slides._sldIdLst  # noqa: SLF001
            ]

            for section_elem in section_list.findall("p14:section", nsmap):
                name = section_elem.get("name", "")
                sld_id_lst = section_elem.find("p14:sldIdLst", nsmap)
                if sld_id_lst is not None:
                    first_id_elem = sld_id_lst.find("p14:sldId", nsmap)
                    if first_id_elem is not None:
                        first_id = first_id_elem.get("id")
                        if first_id is not None:
                            try:
                                start_idx = slide_ids.index(int(first_id))
                            except (ValueError, TypeError):
                                start_idx = len(results)
                            results.append((start_idx, name))

            return results
        except Exception:  # noqa: BLE001
            return []

    # ------------------------------------------------------------------ #
    # XLSX                                                                #
    # ------------------------------------------------------------------ #

    def _detect_xlsx(
        self,
        path: Path,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect worksheets in an Excel workbook.

        Parameters
        ----------
        path:
            Path to the ``.xlsx`` file.

        Returns
        -------
        List[DocumentSection]
        """
        if not _HAS_OPENPYXL:
            raise DetectionError(
                "openpyxl is required for XLSX detection.  "
                "Install it with: pip install openpyxl",
                path=str(path),
            )

        try:
            wb = load_workbook(str(path), read_only=True, data_only=True)
        except Exception as exc:
            raise DetectionError(
                f"Failed to open XLSX file: {exc}", path=str(path)
            ) from exc

        sections: List[DocumentSection] = []
        try:
            for idx, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]
                row_count = ws.max_row or 0
                col_count = ws.max_column or 0
                sections.append(
                    DocumentSection(
                        title=sheet_name,
                        start_index=idx,
                        end_index=idx + 1,
                        level=1,
                        section_type=SectionType.SHEET,
                        detection_method="xlsx_worksheet",
                        confidence=1.0,
                        content_preview=f"{row_count} rows x {col_count} cols",
                        metadata={"row_count": row_count, "col_count": col_count},
                    )
                )
        finally:
            wb.close()

        return sections

    # ------------------------------------------------------------------ #
    # Markdown                                                            #
    # ------------------------------------------------------------------ #

    def _detect_markdown(
        self,
        path: Path,
        heading_level: int = 3,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect sections in a Markdown file by ATX headings.

        YAML front matter (delimited by ``---``) is excluded from detection.
        Headings inside fenced code blocks are ignored.

        Parameters
        ----------
        path:
            Path to the ``.md`` / ``.markdown`` file.
        heading_level:
            Maximum heading level to consider.

        Returns
        -------
        List[DocumentSection]
        """
        text = path.read_text(encoding="utf-8", errors="replace")

        # Strip YAML front matter.
        content = text
        fm_end = 0
        if text.startswith("---"):
            match = re.match(r"^---\r?\n.*?\r?\n---\r?\n?", text, re.DOTALL)
            if match:
                fm_end = match.end()
                content = text[fm_end:]

        # Build a set of character offsets that fall inside fenced code blocks
        # so that we can skip headings occurring there.
        fence_ranges = self._markdown_fence_ranges(content)

        sections: List[DocumentSection] = []
        for m in self._MD_HEADING_RE.finditer(content):
            level = len(m.group(1))
            if level > heading_level:
                continue
            abs_start = fm_end + m.start()
            if self._offset_in_ranges(m.start(), fence_ranges):
                continue
            title = m.group(2).strip()
            sections.append(
                DocumentSection(
                    title=title,
                    start_index=abs_start,
                    end_index=abs_start,  # updated below
                    level=level,
                    section_type=SectionType.HEADING,
                    detection_method="markdown_heading",
                    confidence=1.0,
                    content_preview="",
                )
            )

        # Back-fill end indices.
        total_length = len(text)
        for i, sec in enumerate(sections):
            if i + 1 < len(sections):
                sec.end_index = sections[i + 1].start_index
            else:
                sec.end_index = total_length

        return sections

    # ------------------------------------------------------------------ #
    # LaTeX                                                               #
    # ------------------------------------------------------------------ #

    def _detect_latex(
        self,
        path: Path,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        r"""Detect sections in a LaTeX document by sectioning commands.

        Recognised commands: ``\chapter``, ``\section``, ``\subsection``,
        ``\subsubsection``.

        Parameters
        ----------
        path:
            Path to the ``.tex`` / ``.latex`` file.

        Returns
        -------
        List[DocumentSection]
        """
        text = path.read_text(encoding="utf-8", errors="replace")

        # Remove comments (lines starting with %).
        uncommented = re.sub(r"(?m)(?<!\\)%.*$", "", text)

        sections: List[DocumentSection] = []
        for m in self._LATEX_SECTION_RE.finditer(uncommented):
            cmd = m.group(1)
            title = m.group(2).strip()
            level = self._LATEX_LEVEL_MAP.get(cmd, 2)
            sections.append(
                DocumentSection(
                    title=title,
                    start_index=m.start(),
                    end_index=m.start(),  # updated below
                    level=level,
                    section_type=SectionType.SECTION,
                    detection_method="latex_command",
                    confidence=1.0,
                    content_preview="",
                )
            )

        total_length = len(text)
        for i, sec in enumerate(sections):
            if i + 1 < len(sections):
                sec.end_index = sections[i + 1].start_index
            else:
                sec.end_index = total_length

        return sections

    # ------------------------------------------------------------------ #
    # HTML                                                                #
    # ------------------------------------------------------------------ #

    def _detect_html(
        self,
        path: Path,
        heading_level: int = 3,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect sections in an HTML file by heading tags.

        Parameters
        ----------
        path:
            Path to the ``.html`` / ``.htm`` file.
        heading_level:
            Maximum heading level to consider.

        Returns
        -------
        List[DocumentSection]
        """
        text = path.read_text(encoding="utf-8", errors="replace")

        sections: List[DocumentSection] = []
        for m in self._HTML_HEADING_RE.finditer(text):
            tag = m.group(1).lower()
            level = int(tag[1])
            if level > heading_level:
                continue
            # Strip nested HTML tags from the heading text.
            raw_title = m.group(2)
            title = re.sub(r"<[^>]+>", "", raw_title).strip() or f"{tag} (untitled)"
            sections.append(
                DocumentSection(
                    title=title,
                    start_index=m.start(),
                    end_index=m.start(),
                    level=level,
                    section_type=SectionType.HEADING,
                    detection_method="html_heading",
                    confidence=1.0,
                    content_preview="",
                )
            )

        total_length = len(text)
        for i, sec in enumerate(sections):
            if i + 1 < len(sections):
                sec.end_index = sections[i + 1].start_index
            else:
                sec.end_index = total_length

        return sections

    # ------------------------------------------------------------------ #
    # CSV                                                                 #
    # ------------------------------------------------------------------ #

    def _detect_csv(
        self,
        path: Path,
        rows_per_chunk: int = 1000,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect row-range chunks in a CSV file.

        Parameters
        ----------
        path:
            Path to the ``.csv`` / ``.tsv`` file.
        rows_per_chunk:
            Number of data rows per chunk.

        Returns
        -------
        List[DocumentSection]
        """
        text = path.read_text(encoding="utf-8", errors="replace")
        dialect: Any = csv.excel
        try:
            sample = text[:8192]
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            if path.suffix.lower() == ".tsv":
                dialect = csv.excel_tab

        reader = csv.reader(io.StringIO(text), dialect=dialect)
        rows = list(reader)

        if not rows:
            return []

        # Row 0 is treated as the header.
        total_data_rows = len(rows) - 1
        if total_data_rows <= 0:
            return [
                DocumentSection(
                    title="All rows (header only)",
                    start_index=0,
                    end_index=1,
                    level=1,
                    section_type=SectionType.ROW_RANGE,
                    detection_method="csv_chunk",
                    confidence=1.0,
                )
            ]

        sections: List[DocumentSection] = []
        chunk_idx = 0
        start = 1  # skip header row
        while start <= total_data_rows:
            end = min(start + rows_per_chunk, total_data_rows + 1)
            sections.append(
                DocumentSection(
                    title=f"Rows {start}-{end - 1}",
                    start_index=start,
                    end_index=end,
                    level=1,
                    section_type=SectionType.ROW_RANGE,
                    detection_method="csv_chunk",
                    confidence=1.0,
                    metadata={"chunk_index": chunk_idx},
                )
            )
            chunk_idx += 1
            start = end

        return sections

    # ------------------------------------------------------------------ #
    # JSON                                                                #
    # ------------------------------------------------------------------ #

    def _detect_json(
        self,
        path: Path,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect top-level array elements in a JSON file.

        If the file contains a top-level object instead of an array, a single
        section representing the entire document is returned.

        Parameters
        ----------
        path:
            Path to the ``.json`` file.

        Returns
        -------
        List[DocumentSection]
        """
        text = path.read_text(encoding="utf-8", errors="replace")
        try:
            data = json.loads(text)
        except json.JSONDecodeError as exc:
            raise DetectionError(
                f"Invalid JSON: {exc}", path=str(path)
            ) from exc

        sections: List[DocumentSection] = []

        if isinstance(data, list):
            for idx, item in enumerate(data):
                # Try to derive a meaningful title.
                title = self._json_element_title(item, idx)
                preview = json.dumps(item, ensure_ascii=False)[:120]
                sections.append(
                    DocumentSection(
                        title=title,
                        start_index=idx,
                        end_index=idx + 1,
                        level=1,
                        section_type=SectionType.ARRAY_ELEMENT,
                        detection_method="json_array",
                        confidence=1.0,
                        content_preview=preview,
                    )
                )
        else:
            # Top-level object: treat as a single section.
            sections.append(
                DocumentSection(
                    title="Document",
                    start_index=0,
                    end_index=1,
                    level=1,
                    section_type=SectionType.SECTION,
                    detection_method="json_object",
                    confidence=1.0,
                    content_preview=json.dumps(data, ensure_ascii=False)[:120],
                )
            )

        return sections

    # ------------------------------------------------------------------ #
    # XML                                                                 #
    # ------------------------------------------------------------------ #

    def _detect_xml(
        self,
        path: Path,
        **_kwargs: Any,
    ) -> List[DocumentSection]:
        """Detect top-level child elements in an XML file.

        Parameters
        ----------
        path:
            Path to the ``.xml`` file.

        Returns
        -------
        List[DocumentSection]
        """
        try:
            tree = ElementTree.parse(str(path))  # noqa: S314
        except ElementTree.ParseError as exc:
            raise DetectionError(
                f"Invalid XML: {exc}", path=str(path)
            ) from exc

        root = tree.getroot()
        sections: List[DocumentSection] = []

        for idx, child in enumerate(root):
            # Strip namespace prefix for a cleaner title.
            tag = child.tag
            if "}" in tag:
                tag = tag.split("}", 1)[1]
            text_content = (child.text or "").strip()
            title = tag
            if text_content:
                title = f"{tag}: {text_content[:60]}"
            sections.append(
                DocumentSection(
                    title=title,
                    start_index=idx,
                    end_index=idx + 1,
                    level=1,
                    section_type=SectionType.SECTION,
                    detection_method="xml_child",
                    confidence=1.0,
                    content_preview=ElementTree.tostring(
                        child, encoding="unicode"
                    )[:120],
                )
            )

        return sections

    # ------------------------------------------------------------------ #
    # Document-info enrichment helpers                                    #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _enrich_docx_info(path: Path, info: DocumentInfo) -> None:
        """Populate *info* with DOCX-specific metadata."""
        doc = DocxDocument(str(path))
        word_count = 0
        has_images = False
        has_tables = bool(doc.tables)
        for para in doc.paragraphs:
            word_count += len(para.text.split())
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                has_images = True
                break
        info.word_count = word_count
        info.has_images = has_images
        info.has_tables = has_tables
        core_props = doc.core_properties
        info.metadata.update(
            {
                "author": core_props.author or "",
                "title": core_props.title or "",
                "created": str(core_props.created) if core_props.created else "",
            }
        )

    @staticmethod
    def _enrich_pptx_info(path: Path, info: DocumentInfo) -> None:
        """Populate *info* with PPTX-specific metadata."""
        prs = Presentation(str(path))
        info.page_count = len(prs.slides)
        word_count = 0
        has_images = False
        has_tables = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        word_count += len(para.text.split())
                if shape.shape_type is not None:
                    type_val = int(shape.shape_type)
                    if type_val == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True
                    if type_val == 19:  # MSO_SHAPE_TYPE.TABLE
                        has_tables = True
        info.word_count = word_count
        info.has_images = has_images
        info.has_tables = has_tables

    @staticmethod
    def _enrich_xlsx_info(path: Path, info: DocumentInfo) -> None:
        """Populate *info* with XLSX-specific metadata."""
        wb = load_workbook(str(path), read_only=True, data_only=True)
        try:
            info.page_count = len(wb.sheetnames)
        finally:
            wb.close()

    @staticmethod
    def _enrich_text_info(path: Path, info: DocumentInfo) -> None:
        """Populate *info* with generic text-file metadata."""
        text = path.read_text(encoding="utf-8", errors="replace")
        info.word_count = len(text.split())
        info.has_tables = "<table" in text.lower() or "|" in text

    # ------------------------------------------------------------------ #
    # Internal helpers                                                    #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _markdown_fence_ranges(text: str) -> List[range]:
        """Return character ranges of fenced code blocks in *text*.

        Handles backtick and tilde fences (````` and ``~~~``).
        """
        fence_re = re.compile(r"^(`{3,}|~{3,})[^\S\n]*.*$", re.MULTILINE)
        ranges: List[range] = []
        open_fence: Optional[re.Match[str]] = None
        open_char: Optional[str] = None

        for m in fence_re.finditer(text):
            marker = m.group(1)
            char = marker[0]
            if open_fence is None:
                open_fence = m
                open_char = char
            elif char == open_char and len(marker) >= len(open_fence.group(1)):
                ranges.append(range(open_fence.start(), m.end()))
                open_fence = None
                open_char = None

        # If a fence is never closed, treat everything from it to EOF as fenced.
        if open_fence is not None:
            ranges.append(range(open_fence.start(), len(text)))

        return ranges

    @staticmethod
    def _offset_in_ranges(offset: int, ranges: List[range]) -> bool:
        """Return ``True`` if *offset* falls inside any of *ranges*."""
        for r in ranges:
            if offset in r:
                return True
        return False

    @staticmethod
    def _json_element_title(item: Any, idx: int) -> str:
        """Derive a human-readable title for a JSON array element."""
        if isinstance(item, dict):
            for key in ("title", "name", "id", "label", "key", "heading"):
                if key in item:
                    val = item[key]
                    if isinstance(val, str):
                        return val[:80]
                    return str(val)[:80]
        return f"Element {idx}"
