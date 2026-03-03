"""Document splitting for DOCX, PPTX, XLSX, Markdown, LaTeX, HTML, CSV, JSON, and XML.

The :class:`DocumentSplitter` consumes
:class:`~lazy_splitter.document.models.DocumentSection` segments produced by
:class:`~lazy_splitter.document.detector.DocumentDetector` and writes one
output file per segment.
"""

from __future__ import annotations

import copy
import csv
import io
import json
import logging
import re
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Sequence
from xml.etree import ElementTree

from lazy_splitter.core.base import BaseSplitter
from lazy_splitter.core.exceptions import FileTypeError, SplitError
from lazy_splitter.core.models import DetectionResult
from lazy_splitter.core.utils import ensure_dir, sanitize_filename
from lazy_splitter.document.models import (
    DocumentSection,
    FileType,
    SectionType,
    file_type_from_path,
)

# ---------------------------------------------------------------------------
# Optional heavyweight imports
# ---------------------------------------------------------------------------
try:
    from docx import Document as DocxDocument  # type: ignore[import-untyped]
    from docx.opc.part import Part as DocxPart  # type: ignore[import-untyped]

    _HAS_DOCX = True
except ImportError:  # pragma: no cover
    _HAS_DOCX = False

try:
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Emu  # type: ignore[import-untyped]

    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _HAS_PPTX = False

try:
    from openpyxl import Workbook, load_workbook  # type: ignore[import-untyped]
    from openpyxl.utils import get_column_letter  # type: ignore[import-untyped]

    _HAS_OPENPYXL = True
except ImportError:  # pragma: no cover
    _HAS_OPENPYXL = False


class DocumentSplitter(BaseSplitter):
    """Split documents into smaller files based on detected sections.

    The public entry-point is :meth:`split`, which dispatches to a
    format-specific handler based on the file extension.

    Parameters
    ----------
    logger:
        Optional logger instance.
    """

    # ------------------------------------------------------------------ #
    # Abstract interface implementation                                   #
    # ------------------------------------------------------------------ #

    @property
    def supported_extensions(self) -> List[str]:
        """File extensions this splitter can handle."""
        return [
            ".docx", ".pptx", ".xlsx",
            ".md", ".markdown",
            ".tex", ".latex",
            ".html", ".htm",
            ".csv", ".tsv",
            ".json",
            ".xml",
        ]

    def preview(
        self,
        input_path: Path,
        **kwargs: Any,
    ) -> DetectionResult:
        """Preview what the splitter would produce without writing files.

        Delegates to :class:`~lazy_splitter.document.detector.DocumentDetector`
        and returns the detection result.

        Parameters
        ----------
        input_path:
            Path to the source document.
        **kwargs:
            Strategy-specific options forwarded to the detector.

        Returns
        -------
        DetectionResult
        """
        # Import lazily to avoid circular imports.
        from lazy_splitter.document.detector import DocumentDetector

        detector = DocumentDetector(logger=self.logger)
        return detector.detect(Path(input_path), **kwargs)

    # ------------------------------------------------------------------ #
    # Public API                                                          #
    # ------------------------------------------------------------------ #

    def split(
        self,
        input_path: Path,
        chapters: Sequence[Any],
        **kwargs: Any,
    ) -> List[Path]:
        """Split the document at *input_path* according to *chapters*.

        Parameters
        ----------
        input_path:
            Path to the source document.
        chapters:
            Sequence of :class:`DocumentSection` instances that describe the
            boundaries for splitting.
        **kwargs:
            Format-specific options forwarded to the handler.  Recognised
            keys include:

            * *output_dir* (Path) -- directory to write output files.
              Defaults to a sub-directory next to *input_path* named
              ``<stem>_parts``.
            * *progress_callback* (Callable[[int, int], None]) -- optional
              callable invoked with ``(current_index, total)`` after each
              part is written.

        Returns
        -------
        List[Path]
            Paths to the created output files.

        Raises
        ------
        FileTypeError
            If the file type is unsupported.
        SplitError
            If splitting fails.
        """
        path = Path(input_path)
        if not path.is_file():
            raise SplitError(f"File not found: {path}", path=str(path))

        sections: List[DocumentSection] = list(chapters)  # type: ignore[arg-type]
        if not sections:
            raise SplitError("No sections provided for splitting.", path=str(path))

        output_dir: Optional[Path] = kwargs.pop("output_dir", None)
        progress_callback: Optional[Callable[[int, int], None]] = kwargs.pop(
            "progress_callback", None
        )

        if output_dir is None:
            output_dir = path.parent / f"{path.stem}_parts"
        output_dir = ensure_dir(Path(output_dir))

        file_type = file_type_from_path(path)

        dispatch = {
            FileType.DOCX: self._split_docx,
            FileType.PPTX: self._split_pptx,
            FileType.XLSX: self._split_xlsx,
            FileType.MARKDOWN: self._split_markdown,
            FileType.LATEX: self._split_latex,
            FileType.HTML: self._split_html,
            FileType.CSV: self._split_csv,
            FileType.JSON: self._split_json,
            FileType.XML: self._split_xml,
        }

        handler = dispatch.get(file_type)
        if handler is None:
            raise FileTypeError(
                f"Unsupported file type for splitting: {file_type.value!r}",
                file_type=file_type.value,
            )

        try:
            output_files = handler(
                path,
                sections,
                output_dir,
                progress_callback=progress_callback,
                **kwargs,
            )
        except (SplitError, FileTypeError):
            raise
        except Exception as exc:
            raise SplitError(
                f"Splitting failed for {path.name}: {exc}",
                path=str(path),
                original_error=str(exc),
            ) from exc

        self.logger.info(
            "Split %s into %d file(s) -> %s",
            path.name,
            len(output_files),
            output_dir,
        )

        return output_files

    # ------------------------------------------------------------------ #
    # DOCX                                                                #
    # ------------------------------------------------------------------ #

    def _split_docx(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split a Word document by heading-based sections.

        Each output file is a valid ``.docx`` with styles, images,
        headers and footers copied from the source document.
        """
        if not _HAS_DOCX:
            raise SplitError(
                "python-docx is required for DOCX splitting.  "
                "Install it with: pip install python-docx",
                path=str(path),
            )

        source_doc = DocxDocument(str(path))
        paragraphs = list(source_doc.paragraphs)
        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            new_doc = DocxDocument(str(path))
            # Remove all body content from the clone (keeps styles, headers,
            # footers, and document properties intact).
            body = new_doc.element.body
            for child in list(body):
                body.remove(child)

            # Determine which source paragraphs fall within this section.
            in_section = False
            section_started = False
            char_offset = 0

            for para in paragraphs:
                para_start = char_offset
                para_end = char_offset + len(para.text) + 1

                if para_start >= section.start_index and not section_started:
                    in_section = True
                    section_started = True

                if section_started and para_start >= section.end_index:
                    in_section = False
                    break

                if in_section:
                    # Deep-copy the paragraph XML element and append.
                    new_elem = copy.deepcopy(para._element)  # noqa: SLF001
                    body.append(new_elem)

                char_offset = para_end

            # Copy images that are referenced by the new body.
            self._copy_docx_images(source_doc, new_doc)

            filename = self._make_filename(
                section.title, sec_idx, path.suffix
            )
            out_path = output_dir / filename
            new_doc.save(str(out_path))
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    @staticmethod
    def _copy_docx_images(source_doc: Any, target_doc: Any) -> None:
        """Copy image relationships from *source_doc* to *target_doc*.

        Only images actually referenced in the target body are copied so that
        the output file does not include unreferenced blobs.
        """
        try:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore[import-untyped]

            target_body_xml = target_doc.element.body.xml
            source_part = source_doc.part
            target_part = target_doc.part

            for rel_id, rel in source_part.rels.items():
                if "image" in rel.reltype and rel_id in target_body_xml:
                    # The image blob lives on the related part.
                    try:
                        image_part = rel.target_part
                        target_part.relate_to(image_part, rel.reltype, rel_id)
                    except Exception:  # noqa: BLE001
                        pass
        except Exception:  # noqa: BLE001
            # If the copy fails it is non-fatal -- the DOCX will still open
            # but images may appear as placeholders.
            pass

    # ------------------------------------------------------------------ #
    # PPTX                                                                #
    # ------------------------------------------------------------------ #

    def _split_pptx(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split a PowerPoint presentation by slide-range sections.

        Slide masters and layouts are preserved in each output file.
        """
        if not _HAS_PPTX:
            raise SplitError(
                "python-pptx is required for PPTX splitting.  "
                "Install it with: pip install python-pptx",
                path=str(path),
            )

        source_prs = Presentation(str(path))
        source_slides = list(source_prs.slides)
        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            start = section.start_index
            end = section.end_index

            # Build a new presentation by copying the source and removing
            # slides that do not belong to this section.
            new_prs = Presentation(str(path))
            new_slides = list(new_prs.slides)

            # Determine which slide indices to *keep*.
            keep_indices = set(range(start, end))

            # Remove slides in reverse order to preserve indices.
            for slide_idx in range(len(new_slides) - 1, -1, -1):
                if slide_idx not in keep_indices:
                    self._remove_pptx_slide(new_prs, slide_idx)

            filename = self._make_filename(
                section.title, sec_idx, path.suffix
            )
            out_path = output_dir / filename
            new_prs.save(str(out_path))
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    @staticmethod
    def _remove_pptx_slide(prs: Any, slide_index: int) -> None:
        """Remove the slide at *slide_index* from the presentation XML."""
        try:
            slide_id_list = prs.presentation.sldIdLst
            slide_ids = list(slide_id_list)
            if 0 <= slide_index < len(slide_ids):
                sld_id = slide_ids[slide_index]
                slide_id_list.remove(sld_id)

                # Remove the corresponding slide part relationship.
                r_id = sld_id.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                )
                if r_id is not None:
                    try:
                        prs.part.drop_rel(r_id)
                    except Exception:  # noqa: BLE001
                        pass
        except Exception:  # noqa: BLE001
            pass

    # ------------------------------------------------------------------ #
    # XLSX                                                                #
    # ------------------------------------------------------------------ #

    def _split_xlsx(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split an Excel workbook by worksheets or row ranges.

        Formatting (column widths, number formats, fonts) is preserved.
        """
        if not _HAS_OPENPYXL:
            raise SplitError(
                "openpyxl is required for XLSX splitting.  "
                "Install it with: pip install openpyxl",
                path=str(path),
            )

        source_wb = load_workbook(str(path))
        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            if section.section_type == SectionType.SHEET:
                out_path = self._split_xlsx_by_sheet(
                    source_wb, section, sec_idx, path.suffix, output_dir
                )
            elif section.section_type == SectionType.ROW_RANGE:
                out_path = self._split_xlsx_by_rows(
                    source_wb, section, sec_idx, path.suffix, output_dir
                )
            else:
                # Default: treat as sheet.
                out_path = self._split_xlsx_by_sheet(
                    source_wb, section, sec_idx, path.suffix, output_dir
                )

            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        source_wb.close()
        return output_files

    def _split_xlsx_by_sheet(
        self,
        source_wb: Any,
        section: DocumentSection,
        sec_idx: int,
        suffix: str,
        output_dir: Path,
    ) -> Path:
        """Write a single worksheet into its own workbook."""
        sheet_name = section.title
        if sheet_name not in source_wb.sheetnames:
            raise SplitError(
                f"Worksheet {sheet_name!r} not found in workbook.",
                sheet=sheet_name,
            )

        source_ws = source_wb[sheet_name]
        new_wb = Workbook()
        new_ws = new_wb.active
        new_ws.title = sheet_name  # type: ignore[union-attr]

        # Copy column dimensions.
        for col_letter, dim in source_ws.column_dimensions.items():
            new_ws.column_dimensions[col_letter].width = dim.width  # type: ignore[union-attr]
            new_ws.column_dimensions[col_letter].hidden = dim.hidden  # type: ignore[union-attr]

        # Copy row dimensions.
        for row_num, dim in source_ws.row_dimensions.items():
            new_ws.row_dimensions[row_num].height = dim.height  # type: ignore[index]
            new_ws.row_dimensions[row_num].hidden = dim.hidden  # type: ignore[index]

        # Copy cell values and styles.
        for row in source_ws.iter_rows():
            for cell in row:
                new_cell = new_ws.cell(  # type: ignore[union-attr]
                    row=cell.row, column=cell.column, value=cell.value
                )
                if cell.has_style:
                    new_cell.font = copy.copy(cell.font)
                    new_cell.fill = copy.copy(cell.fill)
                    new_cell.border = copy.copy(cell.border)
                    new_cell.alignment = copy.copy(cell.alignment)
                    new_cell.number_format = cell.number_format
                    new_cell.protection = copy.copy(cell.protection)

        # Copy merged cells.
        for merged_range in source_ws.merged_cells.ranges:
            new_ws.merge_cells(str(merged_range))  # type: ignore[union-attr]

        filename = self._make_filename(section.title, sec_idx, suffix)
        out_path = output_dir / filename
        new_wb.save(str(out_path))
        new_wb.close()
        return out_path

    def _split_xlsx_by_rows(
        self,
        source_wb: Any,
        section: DocumentSection,
        sec_idx: int,
        suffix: str,
        output_dir: Path,
    ) -> Path:
        """Write a row range from the first sheet into its own workbook."""
        source_ws = source_wb.active
        new_wb = Workbook()
        new_ws = new_wb.active

        # Copy column dimensions.
        for col_letter, dim in source_ws.column_dimensions.items():
            new_ws.column_dimensions[col_letter].width = dim.width  # type: ignore[union-attr]

        start_row = section.start_index
        end_row = section.end_index
        dest_row = 1

        for row_num in range(start_row, end_row):
            for col_num in range(1, (source_ws.max_column or 1) + 1):
                src_cell = source_ws.cell(row=row_num + 1, column=col_num)
                new_cell = new_ws.cell(  # type: ignore[union-attr]
                    row=dest_row, column=col_num, value=src_cell.value
                )
                if src_cell.has_style:
                    new_cell.font = copy.copy(src_cell.font)
                    new_cell.fill = copy.copy(src_cell.fill)
                    new_cell.border = copy.copy(src_cell.border)
                    new_cell.alignment = copy.copy(src_cell.alignment)
                    new_cell.number_format = src_cell.number_format
            dest_row += 1

        filename = self._make_filename(section.title, sec_idx, suffix)
        out_path = output_dir / filename
        new_wb.save(str(out_path))
        new_wb.close()
        return out_path

    # ------------------------------------------------------------------ #
    # Markdown                                                            #
    # ------------------------------------------------------------------ #

    def _split_markdown(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split a Markdown file by heading sections.

        YAML front matter is preserved in every output file.  The splitter
        does **not** break inside fenced code blocks -- the entire code block
        is kept with whichever section it falls into.
        """
        text = path.read_text(encoding="utf-8", errors="replace")

        # Extract YAML front matter.
        front_matter = ""
        if text.startswith("---"):
            fm_match = re.match(r"^---\r?\n.*?\r?\n---\r?\n?", text, re.DOTALL)
            if fm_match:
                front_matter = fm_match.group(0)

        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            chunk = text[section.start_index : section.end_index]

            # Prepend front matter if the chunk does not already start with it.
            if front_matter and not chunk.startswith("---"):
                chunk = front_matter + "\n" + chunk

            filename = self._make_filename(section.title, sec_idx, ".md")
            out_path = output_dir / filename
            out_path.write_text(chunk, encoding="utf-8")
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    # ------------------------------------------------------------------ #
    # LaTeX                                                               #
    # ------------------------------------------------------------------ #

    def _split_latex(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        r"""Split a LaTeX file by sectioning commands.

        The preamble (everything before ``\begin{document}``) is preserved in
        every output file, and each file is wrapped with
        ``\begin{document}`` / ``\end{document}``.
        """
        text = path.read_text(encoding="utf-8", errors="replace")

        # Extract preamble.
        preamble = ""
        begin_match = re.search(
            r"\\begin\s*\{document\}", text
        )
        if begin_match:
            preamble = text[: begin_match.end()]

        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            chunk = text[section.start_index : section.end_index].strip()

            # Wrap with preamble and document environment.
            if preamble and not chunk.startswith("\\documentclass"):
                content = preamble + "\n\n" + chunk + "\n\n\\end{document}\n"
            else:
                content = chunk

            filename = self._make_filename(section.title, sec_idx, ".tex")
            out_path = output_dir / filename
            out_path.write_text(content, encoding="utf-8")
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    # ------------------------------------------------------------------ #
    # HTML                                                                #
    # ------------------------------------------------------------------ #

    def _split_html(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split an HTML file by heading tags.

        The ``<head>`` block (including ``<style>`` and ``<link>`` tags) is
        preserved in every output file.
        """
        text = path.read_text(encoding="utf-8", errors="replace")

        # Extract <head> ... </head>.
        head_content = ""
        head_match = re.search(
            r"<head\b[^>]*>(.*?)</head>", text, re.IGNORECASE | re.DOTALL
        )
        if head_match:
            head_content = head_match.group(0)

        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            body_chunk = text[section.start_index : section.end_index].strip()

            # Build a minimal valid HTML document.
            html_doc = (
                "<!DOCTYPE html>\n<html>\n"
                + (head_content + "\n" if head_content else "")
                + "<body>\n"
                + body_chunk
                + "\n</body>\n</html>\n"
            )

            filename = self._make_filename(section.title, sec_idx, ".html")
            out_path = output_dir / filename
            out_path.write_text(html_doc, encoding="utf-8")
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    # ------------------------------------------------------------------ #
    # CSV                                                                 #
    # ------------------------------------------------------------------ #

    def _split_csv(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split a CSV file by row ranges, preserving the header row."""
        text = path.read_text(encoding="utf-8", errors="replace")

        dialect: Any = csv.excel
        try:
            dialect = csv.Sniffer().sniff(text[:8192])
        except csv.Error:
            if path.suffix.lower() == ".tsv":
                dialect = csv.excel_tab

        reader = csv.reader(io.StringIO(text), dialect=dialect)
        rows = list(reader)

        if not rows:
            raise SplitError("CSV file is empty.", path=str(path))

        header = rows[0]
        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            start = section.start_index
            end = section.end_index
            chunk_rows = rows[start:end]

            buf = io.StringIO()
            writer = csv.writer(buf, dialect=dialect)
            writer.writerow(header)
            writer.writerows(chunk_rows)

            filename = self._make_filename(section.title, sec_idx, path.suffix)
            out_path = output_dir / filename
            out_path.write_text(buf.getvalue(), encoding="utf-8")
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    # ------------------------------------------------------------------ #
    # JSON                                                                #
    # ------------------------------------------------------------------ #

    def _split_json(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        indent: int = 2,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split a JSON array file into smaller JSON files.

        Each output file contains the array elements described by the
        corresponding section's index range.  If the original file contains a
        single element, it is written as-is.
        """
        text = path.read_text(encoding="utf-8", errors="replace")
        try:
            data = json.loads(text)
        except json.JSONDecodeError as exc:
            raise SplitError(
                f"Invalid JSON: {exc}", path=str(path)
            ) from exc

        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            if isinstance(data, list):
                start = section.start_index
                end = section.end_index
                chunk = data[start:end]
                # If the section covers exactly one element, write it directly
                # rather than wrapping in a list.
                payload = chunk[0] if len(chunk) == 1 else chunk
            else:
                payload = data

            content = json.dumps(payload, indent=indent, ensure_ascii=False)

            filename = self._make_filename(section.title, sec_idx, ".json")
            out_path = output_dir / filename
            out_path.write_text(content + "\n", encoding="utf-8")
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    # ------------------------------------------------------------------ #
    # XML                                                                 #
    # ------------------------------------------------------------------ #

    def _split_xml(
        self,
        path: Path,
        sections: List[DocumentSection],
        output_dir: Path,
        progress_callback: Optional[Callable[[int, int], None]] = None,
        **_kwargs: Any,
    ) -> List[Path]:
        """Split an XML file by top-level child elements.

        Each output file is a valid XML document whose root element wraps the
        children specified by the section's index range.  The root tag and
        namespace declarations are preserved.
        """
        try:
            tree = ElementTree.parse(str(path))  # noqa: S314
        except ElementTree.ParseError as exc:
            raise SplitError(
                f"Invalid XML: {exc}", path=str(path)
            ) from exc

        root = tree.getroot()
        children = list(root)
        total = len(sections)
        output_files: List[Path] = []

        for sec_idx, section in enumerate(sections):
            start = section.start_index
            end = section.end_index
            selected = children[start:end]

            # Build a new root with the same tag and attributes.
            new_root = ElementTree.Element(root.tag, root.attrib)
            for child in selected:
                new_root.append(copy.deepcopy(child))

            new_tree = ElementTree.ElementTree(new_root)
            filename = self._make_filename(section.title, sec_idx, ".xml")
            out_path = output_dir / filename

            new_tree.write(
                str(out_path), encoding="unicode", xml_declaration=True
            )
            output_files.append(out_path)

            if progress_callback is not None:
                progress_callback(sec_idx + 1, total)

        return output_files

    # ------------------------------------------------------------------ #
    # Shared helpers                                                      #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _make_filename(title: str, index: int, suffix: str) -> str:
        """Generate a safe filename for the *index*-th section.

        Parameters
        ----------
        title:
            Section title (will be sanitised).
        index:
            Zero-based section index.
        suffix:
            File extension **including** the leading dot (e.g. ``".docx"``).

        Returns
        -------
        str
            Filename like ``"01_Introduction.docx"``.
        """
        safe = sanitize_filename(title)
        return f"{index + 1:02d}_{safe}{suffix}"
